"""
Replace speaker notes on slides 9–14 of Final_Report_Slides.pptx with
ready-to-deliver speech.  Slides are identified by title content so the
script keeps working even if the index numbers shift.
"""
from pathlib import Path
from pptx import Presentation

PATH = Path("/Users/bonianhan/Projects/CS699/project/final_present/Final_Report_Slides.pptx")

NOTES = {
    "lagrangian":  # Slide 9 — Constraint convergence
"""Looking at this chart, the green line is PPOLag's training cost over twenty epochs, and the dashed line marks our cost limit, d equals zero point one. Notice the descent. Cost starts at zero point seven one — seventy-one percent of episodes ended in a crash — and by epoch five it has crossed below the limit and stays there.

That smooth descent is the signature of a converged controller. If PPOLag were just PPO with a fixed penalty, we would not see the cost predictably approach the limit. Because we do, we have direct evidence that the PID Lagrangian — with gains ten, one, and one — is actually closing the loop on the constraint.""",

    "did not work":  # Slide 10 — What did NOT work
"""I want to spend a moment on three findings that cut against the paper's narrative, because they explain when the CMDP framework breaks.

First, roundabout. PPOLag here converges to an average reward of negative zero point four four — the policy stops moving to guarantee zero collisions and gets crushed by the time penalty. This is the classic freezing-robot failure mode in safe RL.

Second, intersection. Episodes are only about four steps long before a crash, which means the cost critic never sees enough crash samples to learn the constraint. PPOLag and PPO both crash fifty percent of the time. CMDP is not a free lunch in sparse-cost regimes.

Third, CPO. One hundred percent collisions on merge and highway. This is not a bug. Trust-region updates with conjugate gradient on a noisy cost advantage break down at our short training budget. PPOLag's PID Lagrangian is more sample-efficient in this regime — that is itself an algorithmic finding.""",

    "frontier":  # Slide 11 — Reward vs safety frontier
"""This is the single picture that summarises everything. Each dot is one algorithm on one scenario. The horizontal axis is collision rate — lower is safer. The vertical axis is reward — higher is better. So the top-left corner is what we want.

Notice that PPOLag on merge is the only point sitting clearly in that corner: low collisions, high reward. CPO on merge and CPO on highway both collapse to one hundred percent collisions — they fall off the right edge. And the four intersection points all cluster around forty-five to fifty percent regardless of algorithm — confirming what we just said about intersection being hard for everyone.

Statistically, the merge advantage is significant. Collision p equals zero point zero zero zero. Reward p equals zero point zero zero two. Welch t-test, PPOLag against CPO.""",

    "live demo":  # Slide 12 — Demo merge
"""Let me play the demo. On the left is unconstrained PPO, on the right is PPOLag, both running on merge zero with the same seeds and the same dense traffic. Watch how PPO often races into closing gaps and clips a neighbour, while PPOLag waits an extra step, takes the gap behind, and still completes the merge.

By the end you will see PPO crashed in six of twenty-four episodes, PPOLag in three. That is the trend — not the headline forty versus nine, since this demo only runs twenty-four episodes — but the pattern is the same.

The callout below is worth flagging. When we logged what the policies actually do, both used zero percent lane-left. PPOLag's safety strategy is ninety-eight point five percent slower. It learned to brake hard, not to merge cleverly. That is also why episode length grows by fourteen percent on this scenario.""",

    "all four":  # Slide 13 — 4-env demo
"""This is a twenty-second tour: five seconds per environment, each with its own trained policy. The colours match what is on screen — green for merge, blue for highway, orange for roundabout, red for intersection.

Watch the action distributions on the right. Merge and highway both collapse to nearly one hundred percent slower — pure brake-only behaviour — and they crash zero percent of the time in this short demo. That is the lazy optimum we just talked about.

Roundabout and intersection are different. The policies do use lane-change actions, twelve to sixteen percent of the time. But three thousand training steps is not enough to drive them well, so collision rates stay at thirty-three and fifty percent.

The take-home: the paper's claim is robust on linear-flow tasks where geometry resolves the merge for the agent. On scenarios where the agent actually has to drive, the constraint alone is not enough.""",

    "take-aways":  # Slide 14 — Take-aways
"""Three things to take away from this work.

One — decoupling matters. Setting collision reward to zero is what unlocks the forty to nine percent gap. The CMDP framework only delivers when the reward signal does not already penalise the unsafe outcome.

Two — discrete-action PPOLag often collapses to braking. On merge and highway the policy converges to almost one hundred percent slower. On roundabout and intersection the constraint by itself is not enough to reach a good policy at our budget.

Three — honesty beats hype. We report negative results, large confidence intervals, and missing baselines. We surface them rather than hide them. The headline number is real; the limits are real too.""",
}


def update():
    prs = Presentation(str(PATH))
    print(f"opened {len(prs.slides)} slides")

    for i, slide in enumerate(prs.slides):
        # collect title-ish text
        title = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                tf = shape.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    title = " ".join(r.text for r in tf.paragraphs[0].runs).lower()
                    if title:
                        break
        # match by keyword
        for key, body in NOTES.items():
            if key in title:
                ns = slide.notes_slide
                tf = ns.notes_text_frame
                tf.clear()
                # paragraphs separated by blank line
                first = True
                for chunk in body.split("\n\n"):
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    first = False
                    r = p.add_run()
                    r.text = chunk.strip()
                print(f"  slide {i+1}: rewrote notes for '{key}'")
                break

    prs.save(str(PATH))
    print("saved")


if __name__ == "__main__":
    update()
